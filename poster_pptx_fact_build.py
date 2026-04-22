"""Fact-deck — compliant poster build: copy OOXML, apply text, overlay media only.

Used by :file:`build_poster.py` when ``POSTER_BUILD_MODE=fact`` (default). Mapping must match
``ppt/slides/_rels/slide1.xml.rels`` on the fact file (``DS_Capstone_Poster_FINAL_SCRIPT_COPY.pptx``).

Raster sources are **letterboxed** (uniform scale, white margins) to each part’s native pixel size,
not non-uniformly stretched.
"""

from __future__ import annotations

import io
import os
import shutil
import sys
import zipfile
from pathlib import Path

from pptx import Presentation


def maybe_reexec_with_repo_venv(repo_root: str) -> None:
    """If ``repo_root/.venv`` exists and we are not already using it, restart under that interpreter.

    Stops infinite loops via ``POSTER_NO_VENV_REEXEC=1``. Opt out entirely with the same env var
    before launch. Helps ``python build_poster.py`` pick up a project venv even when the shell
    defaults to conda base.
    """
    if os.environ.get("POSTER_NO_VENV_REEXEC"):
        return
    root = Path(repo_root).resolve()
    for name in ("python", "python3"):
        vpy = root / ".venv" / "bin" / name
        if not vpy.is_file():
            continue
        try:
            if Path(sys.executable).resolve() == vpy.resolve():
                return
        except OSError:
            return
        os.environ["POSTER_NO_VENV_REEXEC"] = "1"
        os.execv(str(vpy), [str(vpy), *sys.argv])
    return

# Same as scripts/build_poster_truth_overlay.MEDIA_REPLACE — fact deck embeds image1..5 on slide1.
MEDIA_ARCNAME_TO_PNG: tuple[tuple[str, str], ...] = (
    ("ppt/media/image2.png", "output/figures/target_balance_v2_ordertime.png"),
    ("ppt/media/image3.png", "output/figures/showcase_model_comparison_heatmap.png"),
    ("ppt/media/image4.png", "poster/diagrams/01_data_and_features.png"),
    ("ppt/media/image5.png", "poster/diagrams/02_validation_and_models.png"),
)


def png_bytes_letterboxed_to_canvas(src: Path, target_w: int, target_h: int) -> bytes:
    """Scale ``src`` uniformly to fit inside ``(target_w, target_h)``, center on white (no stretch).

    Stretching wide Mermaid exports into tall ``image4`` slots was crushing diagram text; oversized
    matplotlib PNGs were then aggressively downscaled. Letterboxing preserves aspect ratio; pair
    this with sources generated near the slot aspect / pixel size.
    """
    from io import BytesIO
    from PIL import Image

    im = Image.open(src)
    if im.mode != "RGBA":
        im = im.convert("RGBA")
    sw, sh = im.size
    scale = min(target_w / sw, target_h / sh)
    nw = max(1, int(round(sw * scale)))
    nh = max(1, int(round(sh * scale)))
    if (nw, nh) != (sw, sh):
        im = im.resize((nw, nh), Image.Resampling.LANCZOS)
    canvas = Image.new("RGBA", (target_w, target_h), (255, 255, 255, 255))
    x = (target_w - nw) // 2
    y = (target_h - nh) // 2
    canvas.paste(im, (x, y), im)
    buf = BytesIO()
    canvas.convert("RGB").save(buf, format="PNG", compress_level=6)
    return buf.getvalue()


def apply_media_overlay_preserve_package(pptx_path: str, repo_root: str) -> int:
    """Replace only listed ``ppt/media/*.png``; preserve zip structure and all other part bytes."""
    from PIL import Image

    repo = Path(repo_root)
    pptx = Path(pptx_path)
    replacements: dict[str, bytes] = {}
    with zipfile.ZipFile(pptx, "r") as zin:
        for arcname, rel_src in MEDIA_ARCNAME_TO_PNG:
            if arcname not in zin.namelist():
                raise FileNotFoundError(
                    f"{pptx.name!r} missing {arcname}; not the expected fact layout."
                )
            old = zin.read(arcname)
            w, h = Image.open(io.BytesIO(old)).size
            src = repo / rel_src
            if not src.is_file():
                raise FileNotFoundError(
                    f"Figure for poster not found: {src}. Regenerate (modeling / mermaid / metrics)."
                )
            replacements[arcname] = png_bytes_letterboxed_to_canvas(src, w, h)

    tmp = pptx.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(pptx, "r") as zin, zipfile.ZipFile(tmp, "w") as zout:
        for info in zin.infolist():
            data = zin.read(info.filename)
            if info.filename in replacements:
                data = replacements[info.filename]
            zout.writestr(info, data)
    tmp.replace(pptx)
    return len(replacements)


def run_poster_from_fact(
    repo_root: str,
    fact_pptx: str,
    out_pptx: str,
    apply_westminster: bool = False,
    apply_deck_text: bool = True,
) -> str:
    """
    1) ``shutil.copy2`` fact → output
    2) Optional ``apply_poster_deck_text`` (requires paragraph/run structure to match
       :data:`poster_deck_text.POSTER_DECK_TEXT`; if not, set ``apply_deck_text=False``)
    3) optional westminster recolor
    4) zip media overlay (model/figure paths → existing image slots)
    """
    from westminster_pptx_brand import apply_westminster_brand_to_slide0

    repo = Path(repo_root).resolve()
    os.chdir(repo)
    fact_p = Path(fact_pptx)
    out_p = Path(out_pptx)
    if not fact_p.is_absolute():
        fact_p = (repo / fact_p).resolve()
    else:
        fact_p = fact_p.resolve()
    if not out_p.is_absolute():
        out_p = (repo / out_p).resolve()
    else:
        out_p = out_p.resolve()
    if not fact_p.is_file():
        raise FileNotFoundError(f"Fact deck not found: {fact_p}")
    if fact_p == out_p:
        raise ValueError("Output path must differ from the fact file (read-only fact).")

    shutil.copy2(fact_p, out_p)
    prs = Presentation(str(out_p))
    text_msg = "text unchanged (apply_deck_text=False)."
    if apply_deck_text:
        from poster_deck_text import apply_poster_deck_text

        try:
            apply_poster_deck_text(prs.slides[0])
            text_msg = "text runs updated from poster_deck_text."
        except Exception as e:
            text_msg = (
                f"text not updated (fact deck does not match poster_deck_text spec: {e}). "
                "Fix the fact export or run with POSTER_APPLY_DECK_TEXT=0; media overlay still applied."
            )
    if apply_westminster:
        apply_westminster_brand_to_slide0(prs.slides[0])
    prs.save(str(out_p))
    n = apply_media_overlay_preserve_package(str(out_p), repo_root)
    return f"{n} media part(s) replaced. {text_msg}"


if __name__ == "__main__":
    pass
