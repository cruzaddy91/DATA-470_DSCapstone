"""Align Westminster (Salt Lake City) template geometry with brand colors in a built .pptx.

`Showcase Templates.pptx` ships with a light lavender/magenta band + W marks (~#A468F2 / #9264CE) that
are not the official **Night** primary. This module retints the top band and logo rasters to match
:mod:`westminster_poster_palette` without re-authoring the template in PowerPoint by hand.
"""

from __future__ import annotations

import io
import colorsys
from typing import cast

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.parts.image import ImagePart

# Digital brand — must match westminster_poster_palette.NIGHT / FLINT
RGB_NIGHT = RGBColor(0x21, 0x15, 0x51)  # PMS 274 C — primary purple
RGB_FLINT = RGBColor(0x10, 0x18, 0x20)  # PMS Black 6 C


def _recolor_westminster_w_rgba(
    data,  # PIL Image.load() PixelAccess
    w: int,
    h: int,
    *,
    night_r: int = 0x21,
    night_g: int = 0x15,
    night_b: int = 0x51,
) -> None:
    """In-place: map lavender / violet W artwork to brand Night, keep light backgrounds."""
    for y in range(h):
        for x in range(w):
            r, g, b, a = data[x, y]
            if a < 28:
                continue
            if r > 250 and g > 250 and b > 250:
                continue
            hsv = colorsys.rgb_to_hsv(r / 255.0, g / 255.0, b / 255.0)
            hue, sat, val = hsv
            # Magenta / violet W — typical in template exports; skip near-white
            if val > 0.97 and sat < 0.04:
                continue
            in_hsv_purple = 0.68 < hue < 0.94 and sat > 0.06
            near_lavender_hex = (r - 0xA4) ** 2 + (g - 0x68) ** 2 + (b - 0xF2) ** 2 < 50_000
            in_dom_template_purple = b > 140 and r > 70 and g < 0xD0 and b > g
            if in_hsv_purple or near_lavender_hex or in_dom_template_purple:
                data[x, y] = (night_r, night_g, night_b, a)


def _png_bytes_recolored_w(blob: bytes) -> bytes:
    from PIL import Image

    im = Image.open(io.BytesIO(blob)).convert("RGBA")
    px = im.load()
    w, h = im.size
    _recolor_westminster_w_rgba(px, w, h)
    out = io.BytesIO()
    im.save(out, format="PNG", compress_level=6)
    return out.getvalue()


def apply_westminster_brand_to_slide0(slide) -> None:
    """Set title bar fill to Night; retint embedded W rasters to Night."""
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if getattr(shape, "name", None) == "Rectangle 1":
            try:
                if shape.fill.type is not None:  # type: ignore[union-attr]
                    shape.fill.solid()  # type: ignore[union-attr]
                    shape.fill.fore_color.rgb = RGB_NIGHT  # type: ignore[union-attr]
            except Exception:
                # Non-solid fills or theme-only; try solid anyway
                try:
                    shape.fill.solid()  # type: ignore[union-attr]
                    shape.fill.fore_color.rgb = RGB_NIGHT  # type: ignore[union-attr]
                except Exception:
                    pass
            break

    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        try:
            r_id = shape._pic.blip_rId  # noqa: SLF001 — python-pptx
        except Exception:
            continue
        if r_id is None:
            continue
        img_part = cast(ImagePart, shape.part.related_part(r_id))
        new_blob = _png_bytes_recolored_w(img_part.blob)
        img_part.blob = new_blob
