#!/usr/bin/env python3
"""Print ``POSTER_DECK_TEXT``-shaped run structure from a .pptx (for pasting into poster_deck_text.py).

Default file resolution: ``tools/poster/exports/<name>`` if present, else ``tools/poster/<name>``.
"""
from __future__ import annotations

import os
import sys
from pathlib import Path
from pprint import pprint

from pptx import Presentation

_POSTER_ROOT = Path(__file__).resolve().parents[1]


def _default_pptx(name: str) -> str:
    exports = _POSTER_ROOT / "exports" / name
    if exports.is_file():
        return str(exports)
    return str(_POSTER_ROOT / name)


def main() -> int:
    os.chdir(_POSTER_ROOT)
    path = os.environ.get("POSTER_PPTX") or _default_pptx("DS_Capstone_Poster_FINAL.pptx")
    prs = Presentation(path)
    slide = prs.slides[0]
    order = [
        "TextBox 11",
        "TextBox 17",
        "TextBox 18",
        "TextBox 19",
        "TextBox 20",
        "TextBox 21",
        "TextBox 22",
        "TextBox 23",
        "TextBox 24",
        "TextBox 25",
        "TextBox 26",
        "TextBox 27",
        "TextBox 28",
        "TextBox 2",
        "TextBox 3",
    ]
    deck: dict[str, list[list[str]]] = {}
    for name in order:
        sh = next((s for s in slide.shapes if s.name == name), None)
        if sh is None or not getattr(sh, "has_text_frame", False):
            continue
        tf = sh.text_frame
        deck[name] = []
        for para in tf.paragraphs:
            runs = [r.text for r in para.runs]
            if runs or (para.text or "").strip():
                deck[name].append(runs)
    print(f"# From {path}\n")
    pprint(deck, width=120, sort_dicts=False)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
