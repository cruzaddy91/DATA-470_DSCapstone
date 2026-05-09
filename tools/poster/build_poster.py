"""Capstone poster build.

**Default (``POSTER_BUILD_MODE=fact`` or unset):** Start from the **fact** deck
(``DS_Capstone_Poster_FINAL_SCRIPT_COPY.pptx``) — copy to the output file, apply
:mod:`poster_deck_text` in place, then replace only ``ppt/media/image2..5.png`` with
model/figure rasters (same pixel dimensions as in the fact file). The fact file is
**never** modified.

**Legacy (``POSTER_BUILD_MODE=showcase_template``):** Rebuild from ``Showcase Templates.pptx``
(layout may differ from the hand fact deck).

**Lightweight fact-only entry:** :file:`render_poster_pptx.py` runs the same fact pipeline without
importing this file’s legacy template layout code (useful for a minimal dependency surface).
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
import os

from poster_deck_text import POSTER_DECK_TEXT as _DECK
from poster_pptx_fact_build import maybe_reexec_with_repo_venv, run_poster_from_fact
from westminster_pptx_brand import RGB_FLINT, RGB_NIGHT, apply_westminster_brand_to_slide0
from poster_layout import (
    HEAD_BODY_GAP_EMU,
    HEAD_BODY_PAIRS,
    MIN_BODY_HEIGHT_EMU,
)

# Legacy Showcase studio file (``POSTER_BUILD_MODE=showcase_template`` only).
REPO_ROOT = str(Path(__file__).resolve().parents[2])
EXPORTS_DIR = os.path.join(REPO_ROOT, "tools", "poster", "exports")
SHOWCASE_TEMPLATE = os.path.join(EXPORTS_DIR, "Showcase Templates.pptx")
# Default build output: never the fact file — see ``run_poster_from_fact``.
OUTPUT = os.environ.get(
    "POSTER_PPTX_OUTPUT",
    os.path.join(EXPORTS_DIR, "DS_Capstone_Poster_FULL_RENDER.pptx"),
)
FIG_DIR = os.path.join(REPO_ROOT, "output", "figures")
# Mermaid exports (`bash tools/poster/scripts/render_poster_diagrams.sh`); keep total raster area ~⅓ of slide
DIAGRAM_DIR = os.path.join(REPO_ROOT, "tools", "poster", "mermaid", "diagrams")

# Single font for entire poster (serif; ensure template/theme does not override in PowerPoint)
FONT = "Times New Roman"
# Type scale — kept in sync with the latest DS_Capstone_Poster_FINAL.pptx (poster is source of truth).
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)  # title bar meta lines on purple band

SZ_TITLE_MAIN = 66  # two-line main title; three runs in OOXML (line1 / space / line2)
SZ_TITLE_BAND = 36  # tagline, name, advisor, program line on purple band
SZ_TITLE_TAGLINE = SZ_TITLE_BAND
SZ_TITLE_NAME = SZ_TITLE_BAND
SZ_TITLE_ADVISOR = SZ_TITLE_BAND
SZ_TITLE_META_LINE = SZ_TITLE_BAND

# Section heading boxes (TextBox 17–22): centered, uniform inset via set_heading()
SZ_SECTION = 48

SZ_BOX1_BODY = 35  # Problem column (TextBox 23)
SZ_BODY = 35  # Data & Scope (TextBox 24)

SZ_MODEL_SUBHEAD = 34  # centered subheads incl. colon in one run (TextBox 25)
SZ_MODEL_BODY = 32

SZ_RESULTS_CAPTION = 40  # Evidence line (TextBox 26)

SZ_TAKEAWAY_BODY = 38  # Key findings (TextBox 27)

SZ_LIMIT_HEAD = 44  # Limitations / boosting / Future Work (TextBox 28)
SZ_LIMIT_BODY = 38
SZ_DATA_SOURCE_LABEL = 24
SZ_DATA_SOURCE_DETAIL = 24

SZ_PIPELINE_ARCH = 40  # legacy name; see TextBox 2 data + repo (measured: 28 pt in FINAL)
SZ_TB2_DATA = 28  # TextBox 2: gray data-source lines

# Layout (914400 EMU = 1 in). Poster slide is 44" × 36".
# Slightly roomier gaps for legibility (avoid cramped rasters and body copy).
GAP_SM = Emu(180000)
GAP_MD = Emu(300000)
# Reserve under TextBox 26 “Evidence” line before the figure stack; keep enough for caption, not at figure expense.
CAPTION_LINE = Emu(520000)
LABEL_ABOVE = Emu(360000)
# Text frame padding (body boxes use slightly tighter horizontal padding to fill width)
MARGIN_TF = Emu(36000)
MARGIN_BODY_TF = Emu(22000)
# Same top/side inset for every section heading box (centers text consistently)
HEADING_MARGIN_TOP = Emu(48000)
HEADING_MARGIN_SIDE = Emu(56000)
HEADING_MARGIN_BOTTOM = Emu(24000)
BODY_LINE_SPACING = 1.14


def _build_from_showcase_template() -> None:
    prs = Presentation(SHOWCASE_TEMPLATE)
    slide = prs.slides[0]  # Slide 1: 3-column, 6 headings, with logos
    apply_westminster_brand_to_slide0(slide)  # Night band + W marks (template uses non-brand lavender)
    
    # ── helpers ──────────────────────────────────────────────────────────
    def find_shape(name):
        for s in slide.shapes:
            if s.name == name:
                return s
        raise KeyError(f"Shape '{name}' not found")
    
    
    def shape_bottom(shape):
        return shape.top + shape.height
    
    
    def nudge_body_below_heading() -> None:
        """Push each body text box down so it starts below its section heading (template gaps are tight)."""
        gap = Emu(HEAD_BODY_GAP_EMU)
        for head_id, body_id in HEAD_BODY_PAIRS:
            h = find_shape(head_id)
            b = find_shape(body_id)
            floor = int(shape_bottom(h) + gap)
            if int(b.top) >= floor:
                continue
            delta = floor - int(b.top)
            new_h = int(b.height) - delta
            if new_h < MIN_BODY_HEIGHT_EMU:
                raise RuntimeError(
                    f"Poster layout: {body_id} would be shorter than {MIN_BODY_HEIGHT_EMU} EMU after nudge "
                    f"(heading {head_id} too tall vs body). Shorten section title copy or template."
                )
            b.top = floor
            b.height = new_h
    
    
    def iter_leaf_shapes(shapes):
        """Walk slide shapes; recurse into groups so nested text boxes get font fixes."""
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from iter_leaf_shapes(shape.shapes)
            else:
                yield shape
    
    
    def apply_text_frame_layout(tf, *, margins: Emu | None = None) -> None:
        """Reduce default padding and anchor text to top so copy fills template boxes."""
        m = margins if margins is not None else MARGIN_TF
        tf.word_wrap = True
        tf.margin_left = m
        tf.margin_right = m
        tf.margin_top = m
        tf.margin_bottom = m
        try:
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        except Exception:
            pass
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            pass
    
    
    def apply_body_text_frame_layout(tf) -> None:
        """Body boxes: tighter margins so paragraphs use full width; top-anchored."""
        apply_text_frame_layout(tf, margins=MARGIN_BODY_TF)
    
    
    def enforce_font_on_all_text(slide, font_name: str) -> None:
        """Set explicit Latin typeface on every run (PowerPoint theme can otherwise show mixed fonts)."""
        for shape in iter_leaf_shapes(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_name
    
    
    def enforce_title_bar_text_colors(slide) -> None:
        """Main title lines must be pure black (#000000); tagline/meta stay white on the band.
    
        Called after ``enforce_font_on_all_text`` so theme or master cannot leave off-black / tx1 tints.
        """
        try:
            tf = find_shape("TextBox 11").text_frame
        except KeyError:
            return
        for pi, para in enumerate(tf.paragraphs):
            target = BLACK if pi == 0 else WHITE
            for run in para.runs:
                run.font.color.rgb = target
    
    
    def set_heading(shape, text, size=None, color=BLACK):
        """Section title: single centered line, uniform inset from top of heading shape (see margins below)."""
        if size is None:
            size = SZ_SECTION
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_top = HEADING_MARGIN_TOP
        tf.margin_bottom = HEADING_MARGIN_BOTTOM
        tf.margin_left = HEADING_MARGIN_SIDE
        tf.margin_right = HEADING_MARGIN_SIDE
        try:
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        except Exception:
            pass
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            pass
        p = tf.paragraphs[0]
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = color
        run.font.name = FONT
        return tf
    
    
    def set_heading_runs(shape, run_texts: list, size=None, color=BLACK):
        """Section title with multiple OOXML runs (same as :func:`set_heading` styling)."""
        if size is None:
            size = SZ_SECTION
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_top = HEADING_MARGIN_TOP
        tf.margin_bottom = HEADING_MARGIN_BOTTOM
        tf.margin_left = HEADING_MARGIN_SIDE
        tf.margin_right = HEADING_MARGIN_SIDE
        try:
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        except Exception:
            pass
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            pass
        p = tf.paragraphs[0]
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.alignment = PP_ALIGN.CENTER
        for t in run_texts:
            run = p.add_run()
            run.text = t
            run.font.size = Pt(size)
            run.font.bold = True
            run.font.color.rgb = color
            run.font.name = FONT
        return tf
    
    
    def set_text(shape, text, size=24, bold=False, color=RGB_FLINT, align=PP_ALIGN.LEFT):
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
        return tf
    
    def add_paragraph(
        tf,
        text,
        size=20,
        bold=False,
        color=RGB_FLINT,
        space_before=Pt(6),
        space_after=Pt(2),
        align=PP_ALIGN.LEFT,
        line_spacing=None,
    ):
        p = tf.add_paragraph()
        p.alignment = align
        p.space_before = space_before
        p.space_after = space_after
        if line_spacing is not None:
            try:
                p.line_spacing = line_spacing
            except Exception:
                pass
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
        return p
    
    BODY_COLOR = RGB_FLINT
    
    
    def _style_run(run, *, size_pt: int, bold: bool, color) -> None:
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color
        run.font.name = FONT
    
    
    def add_prose_paragraph(
        tf,
        text: str,
        size: int,
        *,
        color=BODY_COLOR,
        space_before=Pt(0),
        space_after=Pt(5),
    ) -> None:
        """Body text as prose (no bullet glyph)—reads more like a conference poster than a slide."""
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = space_before
        p.space_after = space_after
        try:
            p.line_spacing = BODY_LINE_SPACING
        except Exception:
            pass
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = FONT
    
    
    def fill_prose_block(tf, lines: list[str], size: int, *, first_space_before=Pt(0)) -> None:
        """First line uses the text frame’s initial paragraph; rest are appended (no leading bullets)."""
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
                p.text = line
                p.alignment = PP_ALIGN.LEFT
                p.space_before = first_space_before
                p.space_after = Pt(4)
                try:
                    p.line_spacing = BODY_LINE_SPACING
                except Exception:
                    pass
                for run in p.runs:
                    run.font.size = Pt(size)
                    run.font.color.rgb = BODY_COLOR
                    run.font.name = FONT
            else:
                add_prose_paragraph(tf, line, size, space_before=Pt(4), space_after=Pt(4))
    
    
    def set_first_paragraph(
        tf,
        text: str,
        *,
        size: int,
        bold: bool = False,
        color=BODY_COLOR,
        space_before=Pt(0),
        space_after=Pt(4),
        align=PP_ALIGN.LEFT,
    ) -> None:
        """After ``text_frame.clear()``, style the single initial paragraph (subhead or opening line)."""
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.space_before = space_before
        p.space_after = space_after
        try:
            p.line_spacing = BODY_LINE_SPACING
        except Exception:
            pass
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = FONT
    
    def add_image_to_slide(img_path, left, top, width=None, height=None):
        pic = slide.shapes.add_picture(
            os.path.join(FIG_DIR, img_path),
            left, top, width, height
        )
        return pic
    
    
    def add_diagram_to_slide(filename: str, left, top, width, height):
        """PNG from ``tools/poster/mermaid/diagrams`` (Mermaid renders)."""
        path = os.path.join(DIAGRAM_DIR, filename)
        if not os.path.isfile(path):
            raise FileNotFoundError(
                f"Diagram not found: {path}. Run: bash scripts/render_poster_diagrams.sh"
            )
        return slide.shapes.add_picture(path, left, top, width=width, height=height)
    
    
    def remove_stale_content_pictures(slide, prs) -> None:
        """Drop raster shapes from the slide body so re-builds do not stack duplicate figures over the template.
    
        Keeps pictures in the top band (university / title logos). Poster figures are placed lower.
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    
        h = int(prs.slide_height)
        keep_top = int(h * 0.17)
        to_remove = []
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            if int(shape.top) < keep_top:
                continue
            to_remove.append(shape)
        for shape in reversed(to_remove):
            sp = shape.element
            sp.getparent().remove(sp)
    
    
    remove_stale_content_pictures(slide, prs)
    
    # ── TITLE BAR ────────────────────────────────────────────────────────
    title_shape = find_shape("TextBox 11")
    tf = title_shape.text_frame
    tf.clear()
    apply_text_frame_layout(tf)
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _tb11 = _DECK["TextBox 11"]
    for part in _tb11[0]:
        run = p.add_run()
        run.text = part
        run.font.size = Pt(SZ_TITLE_MAIN)
        run.font.bold = True
        run.font.color.rgb = BLACK
        run.font.name = FONT
    
    # Tagline / meta: white on purple band
    add_paragraph(
        tf,
        _tb11[1][0],
        size=SZ_TITLE_TAGLINE,
        bold=False,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        space_before=Pt(8),
        space_after=Pt(4),
        line_spacing=BODY_LINE_SPACING,
    )
    add_paragraph(
        tf,
        _tb11[2][0],
        size=SZ_TITLE_NAME,
        bold=False,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        space_before=Pt(6),
        space_after=Pt(4),
    )
    add_paragraph(
        tf,
        _tb11[3][0],
        size=SZ_TITLE_ADVISOR,
        bold=False,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        space_before=Pt(4),
        space_after=Pt(4),
    )
    add_paragraph(
        tf,
        _tb11[4][0],
        size=SZ_TITLE_META_LINE,
        bold=False,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        space_before=Pt(4),
        space_after=Pt(4),
    )
    
    # ── HEADING LABELS (48 pt, centered, black) ──
    for n in (17, 18, 19, 20, 21):
        set_heading(
            find_shape(f"TextBox {n}"), _DECK[f"TextBox {n}"][0][0], color=BLACK
        )
    set_heading_runs(find_shape("TextBox 22"), _DECK["TextBox 22"][0], color=BLACK)
    try:
        set_heading(find_shape("TextBox 3"), _DECK["TextBox 3"][0][0], color=BLACK)
    except KeyError:
        pass
    
    nudge_body_below_heading()
    
    # ── BOX 1: Problem & Motivation (TextBox 23) ────────────────────────
    box1 = find_shape("TextBox 23")
    box1_column_bottom = shape_bottom(box1)
    # Reserve upper portion for prose; place figure in the remainder of this column (no overlap)
    box1.height = int(box1.height * 0.45)
    tf = box1.text_frame
    tf.clear()
    apply_body_text_frame_layout(tf)
    
    fill_prose_block(
        tf, [p[0] for p in _DECK["TextBox 23"]], SZ_BOX1_BODY, first_space_before=Pt(0)
    )
    
    img_left = box1.left + Emu(100000)
    img_top = shape_bottom(box1) + GAP_SM
    img_h = max(Emu(2400000), box1_column_bottom - img_top - GAP_SM)
    img_w = Emu(box1.width - 200000)
    add_image_to_slide("target_balance_v2_ordertime.png", img_left, img_top, width=img_w, height=img_h)
    
    # ── BOX 2: Data & Scope (TextBox 24) ────────────────────────────────
    box2 = find_shape("TextBox 24")
    tf = box2.text_frame
    tf.clear()
    apply_body_text_frame_layout(tf)
    
    # First line: "BigQuery" as its own run (matches deck)
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    p0.space_before = Pt(0)
    p0.space_after = Pt(4)
    try:
        p0.line_spacing = BODY_LINE_SPACING
    except Exception:
        pass
    for part in _DECK["TextBox 24"][0]:
        r = p0.add_run()
        r.text = part
        _style_run(r, size_pt=SZ_BODY, bold=False, color=BODY_COLOR)
    add_prose_paragraph(
        tf,
        _DECK["TextBox 24"][1][0],
        SZ_BODY,
        space_before=Pt(4),
        space_after=Pt(4),
    )
    add_prose_paragraph(
        tf,
        _DECK["TextBox 24"][2][0],
        SZ_BODY,
        space_before=Pt(4),
        space_after=Pt(4),
    )
    
    # ── BOX 3: Modeling Approach (TextBox 25) ────────────────────────────
    box3 = find_shape("TextBox 25")
    tf = box3.text_frame
    tf.clear()
    apply_body_text_frame_layout(tf)
    
    set_first_paragraph(
        tf,
        _DECK["TextBox 25"][0][0],
        size=SZ_MODEL_SUBHEAD,
        bold=True,
        color=RGB_NIGHT,
        space_before=Pt(0),
        space_after=Pt(6),
        align=PP_ALIGN.CENTER,
    )
    add_prose_paragraph(
        tf,
        _DECK["TextBox 25"][1][0],
        SZ_MODEL_BODY,
        space_before=Pt(2),
        space_after=Pt(10),
    )
    add_paragraph(
        tf,
        _DECK["TextBox 25"][2][0],
        size=SZ_MODEL_SUBHEAD,
        bold=True,
        color=RGB_NIGHT,
        space_before=Pt(8),
        space_after=Pt(6),
        line_spacing=BODY_LINE_SPACING,
        align=PP_ALIGN.CENTER,
    )
    add_prose_paragraph(
        tf,
        _DECK["TextBox 25"][3][0],
        SZ_MODEL_BODY,
        space_before=Pt(2),
        space_after=Pt(6),
    )
    add_prose_paragraph(
        tf,
        _DECK["TextBox 25"][4][0],
        SZ_MODEL_BODY,
        space_before=Pt(4),
        space_after=Pt(4),
    )
    
    # ── BOX 4: Results — heatmap + Mermaid diagrams (ROC/PR, drift, feat importance omitted to limit raster area) ──
    box4 = find_shape("TextBox 26")
    tf = box4.text_frame
    tf.clear()
    apply_body_text_frame_layout(tf)
    
    p_cap = tf.paragraphs[0]
    p_cap.alignment = PP_ALIGN.CENTER
    p_cap.space_before = Pt(0)
    p_cap.space_after = Pt(4)
    try:
        p_cap.line_spacing = BODY_LINE_SPACING
    except Exception:
        pass
    for t in _DECK["TextBox 26"][0]:
        r_ev = p_cap.add_run()
        r_ev.text = t
        _style_run(r_ev, size_pt=SZ_RESULTS_CAPTION, bold=True, color=RGB_NIGHT)
    
    stack_top = box4.top + CAPTION_LINE + Emu(120000)
    usable_h = int(shape_bottom(box4) - stack_top - Emu(120000))
    # Slightly favor the Mermaid pair (readability at 44" over tiny flowchart text).
    h_heat = int(usable_h * 0.44)
    g_mid = max(int(usable_h * 0.035), 300000)
    h_diag = usable_h - h_heat - g_mid
    if h_diag < 2_900_000:
        h_heat = max(usable_h - 2_900_000 - g_mid, int(usable_h * 0.40))
        h_diag = usable_h - h_heat - g_mid
    
    heatmap_top = stack_top
    heatmap_left = box4.left + Emu(100000)
    heatmap_w = Emu(box4.width - 200000)
    add_image_to_slide("showcase_model_comparison_heatmap.png", heatmap_left, heatmap_top, heatmap_w, Emu(h_heat))
    
    diag_top = heatmap_top + Emu(h_heat) + Emu(g_mid)
    inset = Emu(50000)
    gutter = Emu(80000)
    inner_w = int(box4.width) - int(inset) * 2 - int(gutter)
    half_w = Emu(inner_w // 2)
    left_x = box4.left + inset
    # Was the data-flow mermaid (01_data_and_features.png). Replaced with the
    # outer-temporal PR curves: a single chart that simultaneously shows lift
    # over the random baseline (dashed line at the 0.89% positive rate),
    # head-to-head LR vs Stack vs LightGBM vs XGBoost, and the recall/precision
    # trade-off across thresholds. Earns its slot vs a generic ETL diagram.
    add_image_to_slide(
        "pr_curves_temporal_v2_ordertime.png",
        left_x,
        diag_top,
        half_w,
        Emu(h_diag),
    )
    # Was the validation methodology mermaid (02_validation_and_models.png).
    # Replaced with a brand-styled deployment readiness summary table that
    # presents the same gate logic as a tabular panel: two headline models,
    # outer-temporal metrics, and the honest pass / fail verdicts (model gates
    # PASS, label maturity FAIL).
    add_image_to_slide(
        "deployment_summary_table_v2_ordertime.png",
        left_x + half_w + gutter,
        diag_top,
        half_w,
        Emu(h_diag),
    )
    
    # ── BOX 5: Key Findings (TextBox 27) — text box shortened; figures below in column ──
    box5 = find_shape("TextBox 27")
    box5_column_bottom = shape_bottom(box5)
    box28 = find_shape("TextBox 28")
    # Shorter “Key Findings” block → taller confusion-matrix strip.
    box5.height = int(box5.height * 0.40)
    tf = box5.text_frame
    tf.clear()
    apply_body_text_frame_layout(tf)
    
    fill_prose_block(
        tf, [p[0] for p in _DECK["TextBox 27"]], SZ_TAKEAWAY_BODY, first_space_before=Pt(0)
    )
    
    cm_top = shape_bottom(box5) + GAP_MD
    fig_floor = min(box5_column_bottom, box28.top - int(GAP_MD))
    avail = max(int(fig_floor - cm_top) - int(GAP_SM), 0)
    # Allow a taller 2-up confusion strip at poster size (row-normalized matrices).
    cm_h = min(avail, 12_000_000)
    
    add_image_to_slide(
        "classification_confusion_matrices_v2_ordertime.png",
        box5.left + Emu(100000),
        cm_top,
        Emu(box5.width - 200000),
        Emu(cm_h),
    )
    
    # ── BOX 6: Limitations & Future Work (TextBox 28) ───────────────────
    box6 = find_shape("TextBox 28")
    tf = box6.text_frame
    tf.clear()
    apply_body_text_frame_layout(tf)
    
    _28 = _DECK["TextBox 28"]
    set_first_paragraph(
        tf,
        _28[0][0],
        size=SZ_LIMIT_HEAD,
        bold=True,
        color=RGB_NIGHT,
        space_before=Pt(0),
        space_after=Pt(6),
        align=PP_ALIGN.CENTER,
    )
    add_prose_paragraph(
        tf,
        _28[1][0],
        SZ_LIMIT_BODY,
        space_before=Pt(2),
        space_after=Pt(10),
    )
    add_paragraph(
        tf,
        _28[2][0],
        size=SZ_LIMIT_HEAD,
        bold=True,
        color=RGB_NIGHT,
        space_before=Pt(4),
        space_after=Pt(6),
        line_spacing=BODY_LINE_SPACING,
        align=PP_ALIGN.CENTER,
    )
    add_prose_paragraph(
        tf,
        _28[3][0],
        SZ_LIMIT_BODY,
        space_before=Pt(2),
        space_after=Pt(10),
    )
    add_paragraph(
        tf,
        _28[4][0],
        size=SZ_LIMIT_HEAD,
        bold=True,
        color=RGB_NIGHT,
        space_before=Pt(4),
        space_after=Pt(6),
        line_spacing=BODY_LINE_SPACING,
        align=PP_ALIGN.CENTER,
    )
    add_prose_paragraph(
        tf,
        _28[5][0],
        SZ_LIMIT_BODY,
        space_before=Pt(2),
        space_after=Pt(10),
    )
    # Data + repo line lives in TextBox 2 in FINAL; do not duplicate here
    
    # ── TextBox 2: data source + reproducibility (matches FINAL; gray body) ──
    _ds_gray = RGBColor(0x88, 0x88, 0x88)
    try:
        _tb2 = find_shape("TextBox 2")
        _tf2 = _tb2.text_frame
        _tf2.clear()
        _tf2.margin_left = Emu(91440)
        _tf2.margin_top = Emu(45720)
        _tf2.margin_right = Emu(91440)
        _tf2.margin_bottom = Emu(45720)
        _tf2.word_wrap = True
        try:
            _tf2.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        except Exception:
            pass
        try:
            _tf2.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            pass
        # Two paragraphs: 3 runs + 1 run (see poster_deck_text / FINAL)
        _d2 = _DECK["TextBox 2"]
        _p0 = _tf2.paragraphs[0]
        _p0.alignment = PP_ALIGN.LEFT
        _p0.space_before = Pt(0)
        _p0.space_after = Pt(2)
        for t in _d2[0]:
            _r = _p0.add_run()
            _r.text = t
            _style_run(_r, size_pt=SZ_TB2_DATA, bold=False, color=_ds_gray)
        _p1 = _tf2.add_paragraph()
        _p1.alignment = PP_ALIGN.LEFT
        _p1.space_before = Pt(2)
        _p1.space_after = Pt(0)
        try:
            _p1.line_spacing = BODY_LINE_SPACING
        except Exception:
            pass
        _r1 = _p1.add_run()
        _r1.text = _d2[1][0]
        _style_run(_r1, size_pt=SZ_TB2_DATA, bold=False, color=_ds_gray)
    except KeyError:
        pass
    
    # ── Remove unused slides (keep only slide 1) ────────────────────────
    # Delete slides 2-8 (reverse order to preserve indices)
    for i in range(len(prs.slides) - 1, 0, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    
    # Theme-safe: re-apply poster font to every text run (labels, bullets, dynamic boxes)
    enforce_font_on_all_text(slide, FONT)
    enforce_title_bar_text_colors(slide)
    
    # ── Save ─────────────────────────────────────────────────────────────
    prs.save(OUTPUT)
    print(f"Poster saved to {OUTPUT}")


if __name__ == "__main__":
    maybe_reexec_with_repo_venv(REPO_ROOT)
    _default_fact = os.path.join(EXPORTS_DIR, "DS_Capstone_Poster_FINAL_SCRIPT_COPY.pptx")
    _default_out = os.path.join(EXPORTS_DIR, "DS_Capstone_Poster_FULL_RENDER.pptx")
    if os.environ.get("POSTER_BUILD_MODE", "fact") == "fact":
        _msg = run_poster_from_fact(
            REPO_ROOT,
            os.environ.get("POSTER_FACT_PPTX", _default_fact),
            os.environ.get("POSTER_PPTX_OUTPUT", _default_out),
            apply_westminster=os.environ.get("POSTER_APPLY_WESTMINSTER_BRAND", "0") == "1",
            apply_deck_text=os.environ.get("POSTER_APPLY_DECK_TEXT", "1") == "1",
        )
        print(_msg)
    else:
        _build_from_showcase_template()
