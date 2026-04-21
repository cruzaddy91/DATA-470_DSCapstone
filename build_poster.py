"""Build the capstone poster from the Westminster template.

Formatting constants and copy are aligned with ``DS_Capstone_Poster_FINAL.pptx`` (the deck is the
source of truth). Re-run after hand-edits and update this file if the artifact drifts.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
import os

from poster_layout import (
    HEAD_BODY_GAP_EMU,
    HEAD_BODY_PAIRS,
    MIN_BODY_HEIGHT_EMU,
)

TEMPLATE = "Showcase Templates.pptx"
OUTPUT = "DS_Capstone_Poster_FINAL.pptx"
FIG_DIR = "output/figures"
# Mermaid exports (`bash scripts/render_poster_diagrams.sh`); keep total raster area ~⅓ of slide
DIAGRAM_DIR = "poster/diagrams"

# Single font for entire poster (serif; ensure template/theme does not override in PowerPoint)
FONT = "Times New Roman"
# Type scale — kept in sync with the latest DS_Capstone_Poster_FINAL.pptx (poster is source of truth).
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)  # title bar meta lines on purple band

SZ_TITLE_MAIN = 56
SZ_TITLE_TAGLINE = 32
SZ_TITLE_NAME = 40
SZ_TITLE_ADVISOR = 40
SZ_TITLE_META_LINE = 32  # school line

# Section heading boxes (TextBox 17–22): centered, uniform inset via set_heading()
SZ_SECTION = 48

SZ_BOX1_BODY = 40  # Problem column (TextBox 23)
SZ_BODY = 32  # Data & Scope (TextBox 24)

SZ_MODEL_SUBHEAD = 32  # centered “Baseline…” / “Validation…” + colon runs (TextBox 25)
SZ_MODEL_BODY = 28

SZ_RESULTS_CAPTION = 40  # Evidence line (TextBox 26)

SZ_TAKEAWAY_BODY = 36  # Key findings (TextBox 27)

SZ_LIMIT_HEAD = 44  # Limitations / boosting / Future Work (TextBox 28)
SZ_LIMIT_BODY = 36
SZ_DATA_SOURCE_LABEL = 24
SZ_DATA_SOURCE_DETAIL = 24

SZ_PIPELINE_ARCH = 40  # TextBox 2 (if present on template)

# Layout (914400 EMU = 1 in). Poster slide is 44" × 36".
GAP_SM = Emu(120000)
GAP_MD = Emu(220000)
# Vertical reserve under Results caption before heatmap (long caption safe at SZ_CAPTION)
CAPTION_LINE = Emu(640000)
LABEL_ABOVE = Emu(360000)
# Text frame padding (body boxes use slightly tighter horizontal padding to fill width)
MARGIN_TF = Emu(36000)
MARGIN_BODY_TF = Emu(22000)
# Same top/side inset for every section heading box (centers text consistently)
HEADING_MARGIN_TOP = Emu(48000)
HEADING_MARGIN_SIDE = Emu(56000)
HEADING_MARGIN_BOTTOM = Emu(24000)
BODY_LINE_SPACING = 1.14

prs = Presentation(TEMPLATE)
slide = prs.slides[0]  # Slide 1: 3-column, 6 headings, with logos

# ── helpers ──────────────────────────────────────────────────────────
def find_shape(name):
    for s in slide.shapes:
        if s.name == name:
            return s
    raise KeyError(f"Shape '{name}' not found")


def shape_bottom(shape):
    return shape.top + shape.height


def nudge_body_below_heading() -> None:
    """Push each body text box down so it starts below its section heading (template gaps are tight)."""
    gap = Emu(HEAD_BODY_GAP_EMU)
    for head_id, body_id in HEAD_BODY_PAIRS:
        h = find_shape(head_id)
        b = find_shape(body_id)
        floor = int(shape_bottom(h) + gap)
        if int(b.top) >= floor:
            continue
        delta = floor - int(b.top)
        new_h = int(b.height) - delta
        if new_h < MIN_BODY_HEIGHT_EMU:
            raise RuntimeError(
                f"Poster layout: {body_id} would be shorter than {MIN_BODY_HEIGHT_EMU} EMU after nudge "
                f"(heading {head_id} too tall vs body). Shorten section title copy or template."
            )
        b.top = floor
        b.height = new_h


def iter_leaf_shapes(shapes):
    """Walk slide shapes; recurse into groups so nested text boxes get font fixes."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_leaf_shapes(shape.shapes)
        else:
            yield shape


def apply_text_frame_layout(tf, *, margins: Emu | None = None) -> None:
    """Reduce default padding and anchor text to top so copy fills template boxes."""
    m = margins if margins is not None else MARGIN_TF
    tf.word_wrap = True
    tf.margin_left = m
    tf.margin_right = m
    tf.margin_top = m
    tf.margin_bottom = m
    try:
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    except Exception:
        pass
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass


def apply_body_text_frame_layout(tf) -> None:
    """Body boxes: tighter margins so paragraphs use full width; top-anchored."""
    apply_text_frame_layout(tf, margins=MARGIN_BODY_TF)


def enforce_font_on_all_text(slide, font_name: str) -> None:
    """Set explicit Latin typeface on every run (PowerPoint theme can otherwise show mixed fonts)."""
    for shape in iter_leaf_shapes(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name


def set_heading(shape, text, size=None, color=BLACK):
    """Section title: single centered line, uniform inset from top of heading shape (see margins below)."""
    if size is None:
        size = SZ_SECTION
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_top = HEADING_MARGIN_TOP
    tf.margin_bottom = HEADING_MARGIN_BOTTOM
    tf.margin_left = HEADING_MARGIN_SIDE
    tf.margin_right = HEADING_MARGIN_SIDE
    try:
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    except Exception:
        pass
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = FONT
    return tf

def set_text(shape, text, size=24, bold=False, color=RGBColor(0x33, 0x33, 0x33), align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    return tf

def add_paragraph(
    tf,
    text,
    size=20,
    bold=False,
    color=RGBColor(0x33, 0x33, 0x33),
    space_before=Pt(6),
    space_after=Pt(2),
    align=PP_ALIGN.LEFT,
    line_spacing=None,
):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    if line_spacing is not None:
        try:
            p.line_spacing = line_spacing
        except Exception:
            pass
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    return p

BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)


def _style_run(run, *, size_pt: int, bold: bool, color) -> None:
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    run.font.name = FONT


def set_first_colon_heading(
    tf,
    label_without_colon: str,
    *,
    size_pt: int,
    space_before=Pt(0),
    space_after=Pt(6),
    align=PP_ALIGN.CENTER,
) -> None:
    """Bold dark-blue heading as label + ':' in separate runs (matches hand-edited PPTX)."""
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    try:
        p.line_spacing = BODY_LINE_SPACING
    except Exception:
        pass
    r1 = p.add_run()
    r1.text = label_without_colon
    _style_run(r1, size_pt=size_pt, bold=True, color=DARK_BLUE)
    r2 = p.add_run()
    r2.text = ":"
    _style_run(r2, size_pt=size_pt, bold=True, color=DARK_BLUE)


def add_colon_heading_paragraph(
    tf,
    label_without_colon: str,
    *,
    size_pt: int,
    space_before=Pt(8),
    space_after=Pt(6),
    align=PP_ALIGN.CENTER,
) -> None:
    """Append a centered colon-split subhead (same run structure as ``set_first_colon_heading``)."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    try:
        p.line_spacing = BODY_LINE_SPACING
    except Exception:
        pass
    r1 = p.add_run()
    r1.text = label_without_colon
    _style_run(r1, size_pt=size_pt, bold=True, color=DARK_BLUE)
    r2 = p.add_run()
    r2.text = ":"
    _style_run(r2, size_pt=size_pt, bold=True, color=DARK_BLUE)


def add_prose_paragraph(
    tf,
    text: str,
    size: int,
    *,
    color=BODY_COLOR,
    space_before=Pt(0),
    space_after=Pt(5),
) -> None:
    """Body text as prose (no bullet glyph)—reads more like a conference poster than a slide."""
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_before = space_before
    p.space_after = space_after
    try:
        p.line_spacing = BODY_LINE_SPACING
    except Exception:
        pass
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = FONT


def fill_prose_block(tf, lines: list[str], size: int, *, first_space_before=Pt(0)) -> None:
    """First line uses the text frame’s initial paragraph; rest are appended (no leading bullets)."""
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = line
            p.alignment = PP_ALIGN.LEFT
            p.space_before = first_space_before
            p.space_after = Pt(4)
            try:
                p.line_spacing = BODY_LINE_SPACING
            except Exception:
                pass
            for run in p.runs:
                run.font.size = Pt(size)
                run.font.color.rgb = BODY_COLOR
                run.font.name = FONT
        else:
            add_prose_paragraph(tf, line, size, space_before=Pt(4), space_after=Pt(4))


def set_first_paragraph(
    tf,
    text: str,
    *,
    size: int,
    bold: bool = False,
    color=BODY_COLOR,
    space_before=Pt(0),
    space_after=Pt(4),
    align=PP_ALIGN.LEFT,
) -> None:
    """After ``text_frame.clear()``, style the single initial paragraph (subhead or opening line)."""
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    try:
        p.line_spacing = BODY_LINE_SPACING
    except Exception:
        pass
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT

def add_image_to_slide(img_path, left, top, width=None, height=None):
    pic = slide.shapes.add_picture(
        os.path.join(FIG_DIR, img_path),
        left, top, width, height
    )
    return pic


def add_diagram_to_slide(filename: str, left, top, width, height):
    """PNG from ``poster/diagrams`` (Mermaid renders)."""
    path = os.path.join(DIAGRAM_DIR, filename)
    if not os.path.isfile(path):
        raise FileNotFoundError(
            f"Diagram not found: {path}. Run: bash scripts/render_poster_diagrams.sh"
        )
    return slide.shapes.add_picture(path, left, top, width=width, height=height)


def remove_stale_content_pictures(slide, prs) -> None:
    """Drop raster shapes from the slide body so re-builds do not stack duplicate figures over the template.

    Keeps pictures in the top band (university / title logos). Poster figures are placed lower.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    h = int(prs.slide_height)
    keep_top = int(h * 0.17)
    to_remove = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if int(shape.top) < keep_top:
            continue
        to_remove.append(shape)
    for shape in reversed(to_remove):
        sp = shape.element
        sp.getparent().remove(sp)


remove_stale_content_pictures(slide, prs)

# ── TITLE BAR ────────────────────────────────────────────────────────
title_shape = find_shape("TextBox 11")
tf = title_shape.text_frame
tf.clear()
apply_text_frame_layout(tf)

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = (
    "Predictive Supply Chain Analytics\n"
    "for Backorder Risk"
)
run.font.size = Pt(SZ_TITLE_MAIN)
run.font.bold = True
run.font.color.rgb = BLACK
run.font.name = FONT

# Tagline / meta: white on purple band (matches current deck)
add_paragraph(
    tf,
    "When simpler models suffice; logistic baselines under temporal validation",
    size=SZ_TITLE_TAGLINE,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.CENTER,
    space_before=Pt(8),
    space_after=Pt(4),
    line_spacing=BODY_LINE_SPACING,
)
add_paragraph(
    tf,
    "Addy Cruz",
    size=SZ_TITLE_NAME,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.CENTER,
    space_before=Pt(6),
    space_after=Pt(4),
)
add_paragraph(
    tf,
    "Advisor: Dr. Liang Jingsai",
    size=SZ_TITLE_ADVISOR,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.CENTER,
    space_before=Pt(4),
    space_after=Pt(4),
)
add_paragraph(
    tf,
    "DATA-470 Capstone  |  Data Science  |  Westminster University",
    size=SZ_TITLE_META_LINE,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.CENTER,
    space_before=Pt(4),
    space_after=Pt(4),
)

# ── HEADING LABELS (48 pt, centered, black — uniform inset via set_heading) ──

headings = {
    "TextBox 17": "Problem: why simplicity",
    "TextBox 18": "Data & Scope",
    "TextBox 19": "Baselines (LR vs. LGBM)",
    "TextBox 20": "Metrics, flow & validation",
    "TextBox 21": "Takeaway: LR vs. boosting",
    "TextBox 22": "Limits & when to boost",
}
for name, label in headings.items():
    set_heading(find_shape(name), label, color=BLACK)

nudge_body_below_heading()

# ── BOX 1: Problem & Motivation (TextBox 23) ────────────────────────
box1 = find_shape("TextBox 23")
box1_column_bottom = shape_bottom(box1)
# Reserve upper portion for prose; place figure in the remainder of this column (no overlap)
box1.height = int(box1.height * 0.54)
tf = box1.text_frame
tf.clear()
apply_body_text_frame_layout(tf)

fill_prose_block(
    tf,
    [
        "Inventory and service decisions need transparent inputs and stable scoring—not only leaderboard accuracy. Imbalance and leakage risk mean complex models can look strong under the wrong split.",
        "This work scores backorder risk at order time with pre-fulfillment predictors only, withholding post-order fields so evaluation matches what operations can actually use.",
        "The outcome is highly imbalanced (~3.4% positives). The poster compares a penalized logistic baseline to gradient boosting under a temporal stress test—the evaluation design matters as much as the algorithm.",
    ],
    SZ_BOX1_BODY,
    first_space_before=Pt(0),
)

img_left = box1.left + Emu(100000)
img_top = shape_bottom(box1) + GAP_SM
img_h = max(Emu(2400000), box1_column_bottom - img_top - GAP_SM)
img_w = Emu(box1.width - 200000)
add_image_to_slide("target_balance_v2_ordertime.png", img_left, img_top, width=img_w, height=img_h)

# ── BOX 2: Data & Scope (TextBox 24) ────────────────────────────────
box2 = find_shape("TextBox 24")
tf = box2.text_frame
tf.clear()
apply_body_text_frame_layout(tf)

# First line: "BigQuery" as its own run (matches deck)
p0 = tf.paragraphs[0]
p0.alignment = PP_ALIGN.LEFT
p0.space_before = Pt(0)
p0.space_after = Pt(4)
try:
    p0.line_spacing = BODY_LINE_SPACING
except Exception:
    pass
for part in (
    "Data come from the SAP Supply Chain ",
    "BigQuery",
    " release (Kaggle): sales, delivery, billing, inventory, purchasing, and master tables, aligned at order-line grain.",
):
    r = p0.add_run()
    r.text = part
    _style_run(r, size_pt=SZ_BODY, bold=False, color=BODY_COLOR)
add_prose_paragraph(
    tf,
    "After labeling and filtering, the analysis uses 31,177 order lines from 52,118 source rows (59.8% coverage), with 3.38% positives (1,054 backorders).",
    SZ_BODY,
    space_before=Pt(4),
    space_after=Pt(4),
)
add_prose_paragraph(
    tf,
    "Thirteen order-time predictors and seven missingness indicators are retained; twenty-three post-order attributes are withheld to enforce a leakage-safe contract. The feature set is tabular and linear-friendly—a strong reason to start with a regularized logistic model before adding boosting.",
    SZ_BODY,
    space_before=Pt(4),
    space_after=Pt(4),
)

# ── BOX 3: Modeling Approach (TextBox 25) ────────────────────────────
box3 = find_shape("TextBox 25")
tf = box3.text_frame
tf.clear()
apply_body_text_frame_layout(tf)

set_first_colon_heading(
    tf,
    "Baseline vs. benchmark",
    size_pt=SZ_MODEL_SUBHEAD,
    space_before=Pt(0),
    space_after=Pt(6),
    align=PP_ALIGN.CENTER,
)
add_prose_paragraph(
    tf,
    "Binary classification at order-line grain (backorder vs. not). Penalized logistic regression is the default, auditable estimator. LightGBM is a complexity benchmark: nonlinear uplift must earn its place against the temporal test, not only on easier splits.",
    SZ_MODEL_BODY,
    space_before=Pt(2),
    space_after=Pt(10),
)
add_colon_heading_paragraph(
    tf,
    "Validation & Metrics",
    size_pt=SZ_MODEL_SUBHEAD,
    space_before=Pt(8),
    space_after=Pt(6),
    align=PP_ALIGN.CENTER,
)
add_prose_paragraph(
    tf,
    "Primary generalization assessment uses a temporal holdout (train on earlier periods, test on later periods). A grouped holdout by sales document provides a secondary diagnostic. The decision threshold is set from five-fold stratified out-of-fold scores on the training period.",
    SZ_MODEL_BODY,
    space_before=Pt(2),
    space_after=Pt(6),
)
add_prose_paragraph(
    tf,
    "PR-AUC and ROC-AUC describe ranking and discrimination; F1, precision, recall, and confusion patterns describe the operating point under rarity—low F1 on the temporal test does not by itself invalidate strong ROC-AUC.",
    SZ_MODEL_BODY,
    space_before=Pt(4),
    space_after=Pt(4),
)

# ── BOX 4: Results — heatmap + Mermaid diagrams (ROC/PR, drift, feat importance omitted to limit raster area) ──
box4 = find_shape("TextBox 26")
tf = box4.text_frame
tf.clear()
apply_body_text_frame_layout(tf)

p_cap = tf.paragraphs[0]
p_cap.alignment = PP_ALIGN.CENTER
p_cap.space_before = Pt(0)
p_cap.space_after = Pt(4)
try:
    p_cap.line_spacing = BODY_LINE_SPACING
except Exception:
    pass
r_ev1 = p_cap.add_run()
r_ev1.text = "Evidence: stress splits, lineage, validation (LR vs. LGBM)"
_style_run(r_ev1, size_pt=SZ_RESULTS_CAPTION, bold=True, color=DARK_BLUE)
r_ev2 = p_cap.add_run()
r_ev2.text = ":"
_style_run(r_ev2, size_pt=SZ_RESULTS_CAPTION, bold=True, color=DARK_BLUE)

stack_top = box4.top + CAPTION_LINE + Emu(80000)
usable_h = int(shape_bottom(box4) - stack_top - Emu(100000))
# Split column: ~half heatmap, ~half pair of diagrams (keeps overall visuals near ~⅓ of slide with left + right figures)
h_heat = int(usable_h * 0.50)
g_mid = max(int(usable_h * 0.03), 260000)
h_diag = usable_h - h_heat - g_mid
if h_diag < 2_600_000:
    h_heat = max(usable_h - 2_600_000 - g_mid, int(usable_h * 0.42))
    h_diag = usable_h - h_heat - g_mid

heatmap_top = stack_top
heatmap_left = box4.left + Emu(100000)
heatmap_w = Emu(box4.width - 200000)
add_image_to_slide("showcase_model_comparison_heatmap.png", heatmap_left, heatmap_top, heatmap_w, Emu(h_heat))

diag_top = heatmap_top + Emu(h_heat) + Emu(g_mid)
inset = Emu(90000)
gutter = Emu(100000)
inner_w = int(box4.width) - int(inset) * 2 - int(gutter)
half_w = Emu(inner_w // 2)
left_x = box4.left + inset
add_diagram_to_slide(
    "01_data_and_features.png",
    left_x,
    diag_top,
    half_w,
    Emu(h_diag),
)
add_diagram_to_slide(
    "02_validation_and_models.png",
    left_x + half_w + gutter,
    diag_top,
    half_w,
    Emu(h_diag),
)

# ── BOX 5: Key Findings (TextBox 27) — text box shortened; figures below in column ──
box5 = find_shape("TextBox 27")
box5_column_bottom = shape_bottom(box5)
box28 = find_shape("TextBox 28")
box5.height = int(box5.height * 0.50)
tf = box5.text_frame
tf.clear()
apply_body_text_frame_layout(tf)

fill_prose_block(
    tf,
    [
        "Temporal holdout is the honest deployment bar; grouped validation is easier when train and test overlap in calendar time—strong grouped metrics should not be mistaken for robustness under drift.",
        "On the temporal holdout, logistic regression achieves higher PR-AUC than LightGBM (0.19 vs. 0.08); both retain strong ROC-AUC (~0.85). Here boosting did not buy better ranking quality on the split that matters most.",
        "An auditable linear baseline aligns with threshold policy and review; LightGBM serves as a nonlinear benchmark, not the default operating choice in this study.",
        "Withholding post-order fields removes leakage and lowers grouped F1 versus a legacy view—evaluation design drives the headline as much as model family. Confusion matrices (row-normalized) show recall at the chosen threshold under extreme imbalance.",
    ],
    SZ_TAKEAWAY_BODY,
    first_space_before=Pt(0),
)

cm_top = shape_bottom(box5) + GAP_MD
fig_floor = min(box5_column_bottom, box28.top - int(GAP_MD))
avail = max(int(fig_floor - cm_top) - int(GAP_SM), 0)
cm_h = min(avail, 9_000_000)  # one panel; avoid runaway height

add_image_to_slide(
    "classification_confusion_matrices_v2_ordertime.png",
    box5.left + Emu(100000),
    cm_top,
    Emu(box5.width - 200000),
    Emu(cm_h),
)

# ── BOX 6: Limitations & Future Work (TextBox 28) ───────────────────
box6 = find_shape("TextBox 28")
tf = box6.text_frame
tf.clear()
apply_body_text_frame_layout(tf)

set_first_colon_heading(
    tf,
    "Limitations",
    size_pt=SZ_LIMIT_HEAD,
    space_before=Pt(0),
    space_after=Pt(6),
    align=PP_ALIGN.CENTER,
)
add_prose_paragraph(
    tf,
    "The temporal test has very few positives (0.89%; n = 58), so threshold metrics are high-variance. Post-order fields are excluded by design—no causal claims. Time-varying positivity and partial labeling complicate fixed thresholds.",
    SZ_LIMIT_BODY,
    space_before=Pt(2),
    space_after=Pt(10),
)
add_colon_heading_paragraph(
    tf,
    "When gradient boosting may still help",
    size_pt=SZ_LIMIT_HEAD,
    space_before=Pt(4),
    space_after=Pt(6),
    align=PP_ALIGN.CENTER,
)
add_prose_paragraph(
    tf,
    "Boosting is worth the added complexity if richer entity histories or interactions produce clear PR-AUC gains on the temporal holdout (not only on grouped splits), or if monitoring cost tradeoffs favor recall. Until then, the simpler model is a defensible default.",
    SZ_LIMIT_BODY,
    space_before=Pt(2),
    space_after=Pt(10),
)
add_colon_heading_paragraph(
    tf,
    "Future Work",
    size_pt=SZ_LIMIT_HEAD,
    space_before=Pt(4),
    space_after=Pt(6),
    align=PP_ALIGN.CENTER,
)
add_prose_paragraph(
    tf,
    "Extend the target to shortfall magnitude where labels allow; add entity histories and demand or inventory context as order-time inputs; deploy rolling temporal validation for production monitoring.",
    SZ_LIMIT_BODY,
    space_before=Pt(2),
    space_after=Pt(10),
)
add_paragraph(
    tf,
    "Data Source",
    size=SZ_DATA_SOURCE_LABEL,
    bold=True,
    color=RGBColor(0x66, 0x66, 0x66),
    space_before=Pt(8),
    space_after=Pt(4),
    line_spacing=BODY_LINE_SPACING,
)
add_paragraph(
    tf,
    "SAP Supply Chain Dataset (BigQuery / Kaggle).",
    size=SZ_DATA_SOURCE_DETAIL,
    color=RGBColor(0x88, 0x88, 0x88),
    space_before=Pt(2),
    space_after=Pt(2),
    line_spacing=BODY_LINE_SPACING,
)
add_paragraph(
    tf,
    "Reproducible ETL and modeling code accompany the project.",
    size=SZ_DATA_SOURCE_DETAIL,
    color=RGBColor(0x88, 0x88, 0x88),
    space_before=Pt(0),
    space_after=Pt(0),
    line_spacing=BODY_LINE_SPACING,
)

# ── Optional: diagram column label (TextBox 2 — present on current template/poster) ──
try:
    _pipe = find_shape("TextBox 2")
    _tfp = _pipe.text_frame
    _tfp.clear()
    _tfp.margin_left = Emu(91440)
    _tfp.margin_top = Emu(45720)
    _tfp.margin_right = Emu(91440)
    _tfp.margin_bottom = Emu(45720)
    _tfp.word_wrap = True
    try:
        _tfp.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    except Exception:
        pass
    try:
        _tfp.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    _pp = _tfp.paragraphs[0]
    _pp.alignment = PP_ALIGN.CENTER
    _pp.space_before = Pt(0)
    _pp.space_after = Pt(5)
    _rr = _pp.add_run()
    _rr.text = "Pipeline Architecture:"
    _style_run(_rr, size_pt=SZ_PIPELINE_ARCH, bold=True, color=DARK_BLUE)
except KeyError:
    pass

# ── Remove unused slides (keep only slide 1) ────────────────────────
# Delete slides 2-8 (reverse order to preserve indices)
for i in range(len(prs.slides) - 1, 0, -1):
    rId = prs.slides._sldIdLst[i].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[i]

# Theme-safe: re-apply poster font to every text run (labels, bullets, dynamic boxes)
enforce_font_on_all_text(slide, FONT)

# ── Save ─────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Poster saved to {OUTPUT}")
