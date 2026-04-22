#!/usr/bin/env python3
"""
Guardrail: section headings must not overlap their body text boxes in the built poster PPTX.

Run after build_poster.py. Uses the same pairs/gap as poster_layout.py.
Default file: DS_Capstone_Poster_FINAL.pptx. Override with POSTER_PPTX_OUTPUT (e.g. script copy).
"""
from __future__ import annotations

import os
import sys
from pathlib import Path

REPO = Path(__file__).resolve().parents[1]
if str(REPO) not in sys.path:
    sys.path.insert(0, str(REPO))

from pptx import Presentation

from poster_layout import (
    HEAD_BODY_PAIRS,
    HEAD_BODY_GAP_EMU,
    LAYOUT_VERIFY_TOLERANCE_EMU,
)

PPTX = REPO / os.environ.get("POSTER_PPTX_OUTPUT", "DS_Capstone_Poster_FINAL.pptx")


def _shape_bottom(shape) -> int:
    return int(shape.top) + int(shape.height)


def main() -> int:
    if not PPTX.exists():
        print(f"Not found: {PPTX}", file=sys.stderr)
        return 1
    prs = Presentation(str(PPTX))
    slide = prs.slides[0]
    by_name = {s.name: s for s in slide.shapes}

    required = {n for pair in HEAD_BODY_PAIRS for n in pair}
    missing = required - by_name.keys()
    if missing:
        print(f"Missing shapes on slide 1: {sorted(missing)}", file=sys.stderr)
        return 2

    min_gap = HEAD_BODY_GAP_EMU - LAYOUT_VERIFY_TOLERANCE_EMU
    errors: list[str] = []
    for head_id, body_id in HEAD_BODY_PAIRS:
        h = by_name[head_id]
        b = by_name[body_id]
        need_top = _shape_bottom(h) + min_gap
        if int(b.top) < need_top:
            errors.append(
                f"{body_id}.top ({int(b.top)}) < {head_id} bottom + gap "
                f"({ _shape_bottom(h) } + {HEAD_BODY_GAP_EMU} EMU, tol {LAYOUT_VERIFY_TOLERANCE_EMU})"
            )

    print(f"File: {PPTX}")
    print(f"Checked {len(HEAD_BODY_PAIRS)} heading/body pairs (min gap {HEAD_BODY_GAP_EMU} EMU).")
    if errors:
        print("Layout guardrail failed:", file=sys.stderr)
        for e in errors:
            print(f"  - {e}", file=sys.stderr)
        return 3
    print("OK: no heading/body vertical collisions (within tolerance).")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
