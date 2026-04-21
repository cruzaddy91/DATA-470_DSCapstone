from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
TEMPLATE_PATH = REPO_ROOT / "Showcase Templates.pptx"
OUTPUT_DIR = REPO_ROOT / "output" / "poster"
OUTPUT_PATH = OUTPUT_DIR / "Addy-Cruz-DS-Capstone-Showcase-Poster.pptx"

FIGURE_TARGET_BALANCE = REPO_ROOT / "output" / "figures" / "target_balance_v2_ordertime.png"
FIGURE_FEATURE_IMPORTANCE = REPO_ROOT / "output" / "figures" / "classification_feature_importance_v2_ordertime.png"
FIGURE_MODEL_COMPARISON = REPO_ROOT / "output" / "figures" / "showcase_model_comparison_heatmap.png"
FIGURE_TEMPORAL_SNAPSHOT = REPO_ROOT / "output" / "figures" / "showcase_temporal_snapshot.png"


PURPLE = RGBColor(84, 60, 137)
BLACK = RGBColor(0, 0, 0)


def keep_only_first_slide(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst
    for idx in range(len(prs.slides) - 1, 0, -1):
        rel_id = slide_id_list[idx].rId
        prs.part.drop_rel(rel_id)
        del slide_id_list[idx]


def set_shape_margins(shape, margin_inches: float = 0.14) -> None:
    margin = Inches(margin_inches)
    shape.text_frame.margin_left = margin
    shape.text_frame.margin_right = margin
    shape.text_frame.margin_top = Inches(0.08)
    shape.text_frame.margin_bottom = Inches(0.08)


def set_heading(shape, text: str) -> None:
    shape.text = ""
    set_shape_margins(shape, 0.06)
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = "Times New Roman"
    font.size = Pt(32)
    font.bold = True
    font.color.rgb = BLACK


def set_body_text(shape, lines: list[str]) -> None:
    shape.text = ""
    set_shape_margins(shape)
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(6)
        paragraph.line_spacing = 1.15
        run = paragraph.add_run()
        run.text = line
        font = run.font
        font.name = "Times New Roman"
        font.size = Pt(26)
        font.color.rgb = BLACK


def set_title_block(shape) -> None:
    shape.text = ""
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0

    title = text_frame.paragraphs[0]
    title.alignment = PP_ALIGN.CENTER
    title_run = title.add_run()
    title_run.text = "Predictive Supply Chain Analytics for Backorder Prevention and Inventory Optimization"
    title_font = title_run.font
    title_font.name = "Times New Roman"
    title_font.size = Pt(48)
    title_font.bold = True
    title_font.color.rgb = PURPLE

    name = text_frame.add_paragraph()
    name.alignment = PP_ALIGN.CENTER
    name.space_before = Pt(2)
    name_run = name.add_run()
    name_run.text = "Addy Cruz"
    name_font = name_run.font
    name_font.name = "Times New Roman"
    name_font.size = Pt(30)
    name_font.bold = True
    name_font.color.rgb = BLACK

    affiliation = text_frame.add_paragraph()
    affiliation.alignment = PP_ALIGN.CENTER
    affiliation.space_before = Pt(2)
    affiliation_run = affiliation.add_run()
    affiliation_run.text = "DATA-470 Capstone | Data Science | Westminster University"
    affiliation_font = affiliation_run.font
    affiliation_font.name = "Times New Roman"
    affiliation_font.size = Pt(22)
    affiliation_font.bold = False
    affiliation_font.color.rgb = BLACK


def add_framed_picture(slide, box_shape, image_path: Path, inset_inches: float = 0.12) -> None:
    inset = Inches(inset_inches)
    slide.shapes.add_picture(
        str(image_path),
        box_shape.left + inset,
        box_shape.top + inset,
        width=box_shape.width - (2 * inset),
        height=box_shape.height - (2 * inset),
    )


def build_poster() -> Path:
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE_PATH}")

    for path in (FIGURE_TARGET_BALANCE, FIGURE_FEATURE_IMPORTANCE, FIGURE_MODEL_COMPARISON, FIGURE_TEMPORAL_SNAPSHOT):
        if not path.exists():
            raise FileNotFoundError(f"Figure not found: {path}")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(TEMPLATE_PATH))
    keep_only_first_slide(prs)
    slide = prs.slides[0]

    title_box = slide.shapes[3]
    heading_shapes = {
        1: slide.shapes[4],
        2: slide.shapes[5],
        3: slide.shapes[6],
        4: slide.shapes[7],
        5: slide.shapes[8],
        6: slide.shapes[9],
    }
    body_shapes = {
        1: slide.shapes[10],
        2: slide.shapes[11],
        3: slide.shapes[12],
        4: slide.shapes[13],
        5: slide.shapes[14],
        6: slide.shapes[15],
    }

    set_title_block(title_box)

    set_heading(heading_shapes[1], "Problem & Goal")
    set_body_text(
        body_shapes[1],
        [
            "Backorders reduce fill rate and delay customers.",
            "Overstock ties up cash and increases waste.",
            "Goal: predict backorder risk at order time from SAP ERP data.",
        ],
    )

    set_heading(heading_shapes[2], "Data & Workflow")
    set_body_text(
        body_shapes[2],
        [
            "20+ SAP tables across sales, delivery, inventory, billing, and purchasing.",
            "ETL: raw tables -> master tables -> BRD metrics -> leakage-safe v2 features.",
            "31,177 labeled order lines; 3.38% positive rate.",
        ],
    )

    set_heading(heading_shapes[3], "Target Distribution")
    set_heading(heading_shapes[4], "Model Comparison Metrics")
    set_body_text(
        body_shapes[4],
        [
            "Production choice is based on the strict temporal split, not the easier grouped split.",
            "Heatmap shows F1 and PR-AUC for the four candidate classifiers across the three validation views.",
            "Grouped scores are stronger because train and test overlap in time.",
            "Recent 24-week performance collapses because that test window has only 14 positives.",
        ],
    )

    set_heading(heading_shapes[5], "Key Drivers")
    set_heading(heading_shapes[6], "Selected Temporal Snapshot")

    add_framed_picture(slide, body_shapes[3], FIGURE_TARGET_BALANCE)
    add_framed_picture(slide, body_shapes[4], FIGURE_MODEL_COMPARISON, inset_inches=0.08)
    add_framed_picture(slide, body_shapes[5], FIGURE_FEATURE_IMPORTANCE)
    add_framed_picture(slide, body_shapes[6], FIGURE_TEMPORAL_SNAPSHOT)

    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


if __name__ == "__main__":
    path = build_poster()
    print(path)
